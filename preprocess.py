#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
次世代地震測報 - PPTX 數據提取與結構化預處理腳本 (preprocess.py)
--------------------------------------------------------------
此腳本自動搜尋當前目錄下的所有 .pptx 檔案，提取每張投影片的：
1. 投影片文字內容 (Slide Text)
2. 講者演講備忘錄 (Speaker Notes)
3. 關聯的影像資源 (Slide Images) 並保存到 gradio_lite_app/assets/
最後編譯出一個結構化的 slides.json 檔案。

腳本支援雙模式：
- 優先使用 `python-pptx` 函式庫以獲得最佳相容性與精準解析。
- 若無安裝，會自動切換為 Python 內置 `zipfile` 與 XML 解析的免依賴免安裝模式。
"""

import os
import re
import json
import shutil
import zipfile
import xml.etree.ElementTree as ET

# 目標目錄配置
CURRENT_DIR = os.path.dirname(os.path.abspath(__file__)) if __file__ else os.getcwd()
DEST_DIR = os.path.join(CURRENT_DIR, "gradio_lite_app")
ASSETS_DIR = os.path.join(DEST_DIR, "assets")

# 確保輸出目錄存在
os.makedirs(ASSETS_DIR, exist_ok=True)

def sanitize_filename(name):
    # 清理檔名，只保留英數底線與中文
    return re.sub(r'[^\w\u4e00-\u9fa5]', '_', name)

# ================= 模式一：專業 python-pptx 解析 =================
def process_with_pptx(pptx_path, doc_title, sanitized_doc):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    
    prs = Presentation(pptx_path)
    slides_data = []
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        slide_id = f"slide{slide_num}"
        
        # 1. 提取文字
        text_runs = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text.strip())
        slide_text = " ".join([t for t in text_runs if t])
        
        # 2. 提取講者備忘錄
        notes_text = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            
        # 3. 提取投影片圖片
        slide_images = []
        image_count = 1
        for shape in slide.shapes:
            # 處理一般圖片形狀
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image = shape.image
                ext = image.ext
                # 檔名格式: doc_slideN_imageM.ext
                img_filename = f"{sanitized_doc}_slide{slide_num}_image{image_count}.{ext}"
                img_path = os.path.join(ASSETS_DIR, img_filename)
                
                with open(img_path, "wb") as f:
                    f.write(image.blob)
                
                slide_images.append(f"assets/{img_filename}")
                image_count += 1
                
            # 處理群組形狀中的圖片
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for sub_shape in shape.shapes:
                    if sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = sub_shape.image
                        ext = image.ext
                        img_filename = f"{sanitized_doc}_slide{slide_num}_image{image_count}.{ext}"
                        img_path = os.path.join(ASSETS_DIR, img_filename)
                        
                        with open(img_path, "wb") as f:
                            f.write(image.blob)
                        
                        slide_images.append(f"assets/{img_filename}")
                        image_count += 1
                        
        slides_data.append({
            "doc_title": doc_title,
            "slide_id": slide_id,
            "slide_num": slide_num,
            "content": slide_text,
            "notes": notes_text,
            "images": slide_images
        })
        
    return slides_data

# ================= 模式二：免依賴內置 Zip 讀取與 XML 解析 =================
def process_with_zip(pptx_path, doc_title, sanitized_doc):
    slides_data = []
    
    with zipfile.ZipFile(pptx_path, 'r') as zip_ref:
        # 尋找所有投影片 XML
        slide_entries = [f for f in zip_ref.namelist() if f.startswith("ppt/slides/slide") and f.endswith(".xml")]
        # 排序投影片，確保順序正確 slide1, slide2, slide3...
        slide_entries.sort(key=lambda x: int(re.search(r'slide(\d+)\.xml', x).group(1)))
        
        # 讀取備忘錄
        notes_dict = {}
        notes_entries = [f for f in zip_ref.namelist() if f.startswith("ppt/notesSlides/notesSlide") and f.endswith(".xml")]
        for n_entry in notes_entries:
            n_num = int(re.search(r'notesSlide(\d+)\.xml', n_entry).group(1))
            try:
                xml_data = zip_ref.read(n_entry)
                root = ET.fromstring(xml_data)
                # DrawingML namespace
                a_namespace = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                text_runs = []
                for elem in root.iter(f"{a_namespace}t"):
                    if elem.text:
                        text_runs.append(elem.text.strip())
                notes_dict[n_num] = " ".join([t for t in text_runs if t])
            except Exception as e:
                pass
                
        # 逐頁解析
        for entry in slide_entries:
            slide_num = int(re.search(r'slide(\d+)\.xml', entry).group(1))
            slide_id = f"slide{slide_num}"
            
            # 1. 提取投影片內文字
            xml_data = zip_ref.read(entry)
            root = ET.fromstring(xml_data)
            a_namespace = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            text_runs = []
            for elem in root.iter(f"{a_namespace}t"):
                if elem.text:
                    text_runs.append(elem.text.strip())
            slide_text = " ".join([t for t in text_runs if t])
            
            # 2. 對照備忘錄
            slide_notes = notes_dict.get(slide_num, "")
            
            # 3. 提取圖片 (藉由解析與 slideN.xml 關聯的關係檔 .xml.rels)
            slide_images = []
            rel_entry = f"ppt/slides/_rels/slide{slide_num}.xml.rels"
            
            if rel_entry in zip_ref.namelist():
                rel_data = zip_ref.read(rel_entry)
                # 使用正則搜尋 Target="../media/image1.png"
                img_targets = re.findall(b'Target="\\.\\./media/([^"]+)"', rel_data)
                
                for byte_target in img_targets:
                    img_name = byte_target.decode('utf-8')
                    zip_media_path = f"ppt/media/{img_name}"
                    
                    if zip_media_path in zip_ref.namelist():
                        out_img_name = f"{sanitized_doc}_slide{slide_num}_{img_name}"
                        out_img_path = os.path.join(ASSETS_DIR, out_img_name)
                        
                        # 實體解壓圖片並保存
                        with zip_ref.open(zip_media_path) as source, open(out_img_path, "wb") as target:
                            shutil.copyfileobj(source, target)
                            
                        slide_images.append(f"assets/{out_img_name}")
                        
            slides_data.append({
                "doc_title": doc_title,
                "slide_id": slide_id,
                "slide_num": slide_num,
                "content": slide_text,
                "notes": slide_notes,
                "images": slide_images
            })
            
    return slides_data

# ================= 執行主控台 =================
def main():
    print("🚀 啟動次世代地震測報 PPTX 數據提取程序...")
    
    # 搜尋當前目錄下的 .pptx 檔案
    pptx_files = [f for f in os.listdir(CURRENT_DIR) if f.lower().endswith(".pptx")]
    
    if not pptx_files:
        print("⚠️ 錯誤：在當前目錄中未找到任何 .pptx 檔案。")
        print(f"請確認簡報檔案是否位於：{CURRENT_DIR}")
        return
        
    # 檢測是否有安裝 python-pptx 庫
    has_pptx_lib = False
    try:
        import pptx
        has_pptx_lib = True
        print("💡 檢測到系統已安裝 python-pptx，將採用高級解析模式。")
    except ImportError:
        print("💡 未檢測到 python-pptx，自動切換至內置無依賴 XML-Zip 提取模式。")
        
    all_slides_data = []
    
    for filename in pptx_files:
        pptx_path = os.path.join(CURRENT_DIR, filename)
        doc_title = os.path.splitext(filename)[0]
        sanitized_doc = sanitize_filename(doc_title)
        
        print(f"📄 正在解析：{filename}...")
        
        try:
            if has_pptx_lib:
                doc_data = process_with_pptx(pptx_path, doc_title, sanitized_doc)
            else:
                doc_data = process_with_zip(pptx_path, doc_title, sanitized_doc)
                
            all_slides_data.extend(doc_data)
            print(f"   └─ ✅ 成功解析 {len(doc_data)} 頁投影片。")
        except Exception as e:
            print(f"   └─ ❌ 解析失敗，錯誤原因：{e}")
            
    # 輸出結構化 slides.json 檔案
    json_output_path = os.path.join(DEST_DIR, "slides.json")
    with open(json_output_path, "w", encoding="utf-8") as json_file:
        json.dump(all_slides_data, json_file, ensure_ascii=False, indent=4)
        
    print("\n==========================================")
    print("🎉 資料提煉完成！")
    print(f"📁 圖片資源已提取至：{ASSETS_DIR}")
    print(f"📄 結構化知識庫已生成：{json_output_path}")
    print(f"✨ 總計提煉了 {len(all_slides_data)} 頁投影片內容。")
    print("==========================================")

if __name__ == "__main__":
    main()
